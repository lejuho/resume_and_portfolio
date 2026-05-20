"""Unit tests for preset loading, merging, saving, and CLI integration."""

from __future__ import annotations

import shutil
import textwrap
from pathlib import Path

import pytest
from typer.testing import CliRunner

from scripts.pcli import app
from scripts.preset import PresetError, load_preset
from scripts.select import filter_cards

runner = CliRunner()

# ─── Fixtures ──────────────────────────────────────────────────────────────

SAMPLE_MDX = textwrap.dedent("""\
    ---
    id: sample-card
    title: Sample Card
    type: hackathon
    period:
      start: 2026-05-01
    status: live
    summary: "A valid summary for testing the CLI commands."
    summary_ko: "테스트용 유효한 한국어 요약문입니다."
    narrative: "Long enough narrative to pass live status check. Over 100 chars needed here yes."
    evidence:
      - type: repo
        url: https://github.com/example/repo
    tags:
      domain: [web3]
      skill: [solidity]
      audience: []
    ---
""")

NO_KO_MDX = textwrap.dedent("""\
    ---
    id: no-ko-card
    title: English Only Card
    type: project
    period:
      start: 2026-03-01
    status: live
    summary: "No Korean summary here."
    narrative: "Long enough narrative to pass live status check. Over 100 chars needed here yes."
    evidence:
      - type: repo
        url: https://github.com/example/repo
    ---
""")


@pytest.fixture()
def preset_repo(tmp_path):
    cards_dir = tmp_path / "cards"
    cards_dir.mkdir()
    (cards_dir / "2026-05-sample-card.mdx").write_text(SAMPLE_MDX, encoding="utf-8")
    (cards_dir / "2026-03-no-ko-card.mdx").write_text(NO_KO_MDX, encoding="utf-8")

    profile = textwrap.dedent("""\
        basics:
          name: Test User
          label: Test Engineer
          email: test@example.com
          summary_en: "Test summary."
        education: []
    """)
    (tmp_path / "profile.example.yaml").write_text(profile, encoding="utf-8")
    (tmp_path / ".build").mkdir()
    (tmp_path / "output").mkdir()
    (tmp_path / "presets").mkdir()
    (tmp_path / "templates" / "portfolio").mkdir(parents=True)
    template_src = Path(__file__).parents[1] / "templates" / "portfolio" / "default.py"
    shutil.copy(template_src, tmp_path / "templates" / "portfolio" / "default.py")

    import scripts.pcli as pcli_mod

    pcli_mod.REPO_ROOT = tmp_path
    yield tmp_path
    # restore
    import importlib

    importlib.reload(pcli_mod)


# ─── Preset load ───────────────────────────────────────────────────────────


def test_load_preset_ok(preset_repo):
    p_yaml = textwrap.dedent("""\
        target: portfolio
        layout: one-per-card
        max_items: 5
        filters:
          tags: web3
    """)
    (preset_repo / "presets" / "test.yaml").write_text(p_yaml, encoding="utf-8")
    p = load_preset("test", preset_repo / "presets")
    assert p.target == "portfolio"
    assert p.layout == "one-per-card"
    assert p.max_items == 5
    assert p.filters.tags == "web3"


def test_load_preset_missing(preset_repo):
    with pytest.raises(PresetError, match="not found"):
        load_preset("nonexistent", preset_repo / "presets")


def test_load_preset_missing_target(preset_repo):
    (preset_repo / "presets" / "bad.yaml").write_text("layout: one-per-card\n", encoding="utf-8")
    with pytest.raises(PresetError, match="target"):
        load_preset("bad", preset_repo / "presets")


def test_load_preset_include_exclude(preset_repo):
    p_yaml = textwrap.dedent("""\
        target: portfolio
        include_cards:
          - card-a
          - card-b
        exclude_cards:
          - card-c
    """)
    (preset_repo / "presets" / "inc.yaml").write_text(p_yaml, encoding="utf-8")
    p = load_preset("inc", preset_repo / "presets")
    assert p.include_cards == ["card-a", "card-b"]
    assert p.exclude_cards == ["card-c"]


# ─── select.py exclude_ids ─────────────────────────────────────────────────


def test_exclude_ids_removes_from_result():
    from datetime import date

    from scripts.card import Card

    def _card(id_, type_="project"):
        return Card.model_validate(
            {
                "id": id_,
                "title": id_,
                "type": type_,
                "period": {"start": date(2026, 1, 1)},
                "status": "live",
                "summary": "ok",
                "evidence": [{"type": "repo", "url": "https://example.com"}],
            }
        )

    cards = [_card("a"), _card("b"), _card("c")]
    result = filter_cards(cards, status="live", exclude_ids=["b"])
    assert [c.id for c in result] == ["a", "c"]


def test_exclude_ids_with_explicit_ids():
    from datetime import date

    from scripts.card import Card

    def _card(id_):
        return Card.model_validate(
            {
                "id": id_,
                "title": id_,
                "type": "project",
                "period": {"start": date(2026, 1, 1)},
                "status": "live",
                "summary": "ok",
                "evidence": [{"type": "repo", "url": "https://example.com"}],
            }
        )

    cards = [_card("a"), _card("b"), _card("c")]
    # explicit_ids=["a","b"] but exclude_ids=["b"] → only "a"
    result = filter_cards(cards, explicit_ids=["a", "b"], exclude_ids=["b"])
    assert [c.id for c in result] == ["a"]


# ─── preset save ───────────────────────────────────────────────────────────


def test_preset_save_roundtrip(preset_repo, monkeypatch):
    import scripts.pcli as pcli_mod

    monkeypatch.setattr(pcli_mod, "REPO_ROOT", preset_repo)
    # First build a portfolio to populate cache
    out = preset_repo / "output" / "p.pptx"
    result = runner.invoke(app, ["build", "portfolio", "--out", str(out)])
    assert result.exit_code == 0

    # Save to a preset
    result2 = runner.invoke(app, ["preset", "save", "my-preset"])
    assert result2.exit_code == 0
    saved = preset_repo / "presets" / "my-preset.yaml"
    assert saved.exists()

    # Load it and check it has target
    p = load_preset("my-preset", preset_repo / "presets")
    assert p.target == "portfolio"


def test_preset_save_no_cache(preset_repo, monkeypatch):
    import scripts.pcli as pcli_mod

    monkeypatch.setattr(pcli_mod, "REPO_ROOT", preset_repo)
    # No build has been done — cache is empty
    result = runner.invoke(app, ["preset", "save", "empty"])
    assert result.exit_code != 0


# ─── bok summary_ko enforcement ────────────────────────────────────────────


def test_bok_fails_missing_summary_ko(preset_repo, monkeypatch):
    """bok template should fail before Typst when a card lacks summary_ko (non-dry-run)."""
    import scripts.pcli as pcli_mod

    monkeypatch.setattr(pcli_mod, "REPO_ROOT", preset_repo)
    # no-ko-card has no summary_ko; no --dry-run so validation runs before Typst check
    result = runner.invoke(app, ["build", "resume", "--template", "bok", "--cards", "no-ko-card"])
    assert result.exit_code != 0
    assert "summary_ko" in result.output or "summary_ko" in (result.stderr or "")


def test_bok_passes_when_all_have_summary_ko(preset_repo, monkeypatch):
    """bok dry-run should pass when all selected cards have summary_ko."""
    import scripts.pcli as pcli_mod

    monkeypatch.setattr(pcli_mod, "REPO_ROOT", preset_repo)
    result = runner.invoke(
        app, ["build", "resume", "--template", "bok", "--cards", "sample-card", "--dry-run"]
    )
    assert result.exit_code == 0
    assert "lang=ko" in result.output or "template=bok" in result.output


# ─── portfolio layout variants ─────────────────────────────────────────────


def test_build_portfolio_grouped_by_type_dry_run(preset_repo, monkeypatch):
    import scripts.pcli as pcli_mod

    monkeypatch.setattr(pcli_mod, "REPO_ROOT", preset_repo)
    result = runner.invoke(app, ["build", "portfolio", "--layout", "grouped-by-type", "--dry-run"])
    assert result.exit_code == 0
    assert "grouped-by-type" in result.output


def test_build_portfolio_timeline_dry_run(preset_repo, monkeypatch):
    import scripts.pcli as pcli_mod

    monkeypatch.setattr(pcli_mod, "REPO_ROOT", preset_repo)
    result = runner.invoke(app, ["build", "portfolio", "--layout", "timeline", "--dry-run"])
    assert result.exit_code == 0
    assert "timeline" in result.output


def test_build_portfolio_grouped_by_type_pptx(preset_repo, monkeypatch):
    """grouped-by-type PPTX should contain a type divider slide."""
    import scripts.pcli as pcli_mod

    monkeypatch.setattr(pcli_mod, "REPO_ROOT", preset_repo)
    out = preset_repo / "output" / "grouped.pptx"
    result = runner.invoke(
        app, ["build", "portfolio", "--layout", "grouped-by-type", "--out", str(out)]
    )
    assert result.exit_code == 0
    assert out.exists()

    from pptx import Presentation

    prs = Presentation(str(out))
    all_text = "\n".join(
        shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")
    )
    # Divider slide should contain the card type in uppercase
    assert "HACKATHON" in all_text or "PROJECT" in all_text


def test_build_portfolio_timeline_pptx(preset_repo, monkeypatch):
    """timeline PPTX should contain the timeline overview text."""
    import scripts.pcli as pcli_mod

    monkeypatch.setattr(pcli_mod, "REPO_ROOT", preset_repo)
    out = preset_repo / "output" / "timeline.pptx"
    result = runner.invoke(app, ["build", "portfolio", "--layout", "timeline", "--out", str(out)])
    assert result.exit_code == 0
    assert out.exists()

    from pptx import Presentation

    prs = Presentation(str(out))
    all_text = "\n".join(
        shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")
    )
    assert "Timeline" in all_text


# ─── preset run ────────────────────────────────────────────────────────────


def test_preset_run_portfolio(preset_repo, monkeypatch):
    import scripts.pcli as pcli_mod

    monkeypatch.setattr(pcli_mod, "REPO_ROOT", preset_repo)
    p_yaml = textwrap.dedent("""\
        target: portfolio
        layout: one-per-card
        max_items: 5
    """)
    (preset_repo / "presets" / "smoke.yaml").write_text(p_yaml, encoding="utf-8")
    out = preset_repo / "output" / "smoke.pptx"
    result = runner.invoke(app, ["preset", "run", "smoke", "--out", str(out)])
    assert result.exit_code == 0
    assert out.exists()


def test_preset_run_missing(preset_repo, monkeypatch):
    import scripts.pcli as pcli_mod

    monkeypatch.setattr(pcli_mod, "REPO_ROOT", preset_repo)
    result = runner.invoke(app, ["preset", "run", "doesnotexist"])
    assert result.exit_code != 0
