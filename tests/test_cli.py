"""CLI integration tests using Typer's CliRunner."""

from __future__ import annotations

import shutil
import textwrap
from pathlib import Path

import pytest
from typer.testing import CliRunner

from scripts.pcli import app

runner = CliRunner()


# ─── Fixtures ──────────────────────────────────────────────────────────────

SAMPLE_MDX = textwrap.dedent("""\
    ---
    id: sample-card
    title: Sample Card
    type: hackathon
    period:
      start: 2026-05-01
    status: live
    summary: "A valid summary for testing the CLI commands."
    narrative: "Long enough narrative to pass live status check. Over 100 chars needed here yes."
    evidence:
      - type: repo
        url: https://github.com/example/repo
    ---
""")


@pytest.fixture()
def card_repo(tmp_path, monkeypatch):
    """Monkeypatch REPO_ROOT and create a minimal card."""
    cards_dir = tmp_path / "cards"
    cards_dir.mkdir()
    (cards_dir / "2026-05-sample-card.mdx").write_text(SAMPLE_MDX, encoding="utf-8")
    # profile
    profile = textwrap.dedent("""\
        basics:
          name: Test User
          label: Test Engineer
          email: test@example.com
          summary_en: "Test summary."
        education: []
    """)
    (tmp_path / "profile.example.yaml").write_text(profile, encoding="utf-8")
    # .build and output dirs
    (tmp_path / ".build").mkdir()
    (tmp_path / "output").mkdir()
    (tmp_path / "templates" / "portfolio").mkdir(parents=True)
    template_src = Path(__file__).parents[1] / "templates" / "portfolio" / "default.py"
    shutil.copy(template_src, tmp_path / "templates" / "portfolio" / "default.py")

    import scripts.pcli as pcli_mod
    import scripts.render_resume as rr_mod

    monkeypatch.setattr(pcli_mod, "REPO_ROOT", tmp_path)
    monkeypatch.setattr(rr_mod, "console", rr_mod.console)  # no-op, keep reference
    return tmp_path


# ─── help ──────────────────────────────────────────────────────────────────


def test_help():
    result = runner.invoke(app, ["--help"])
    assert result.exit_code == 0
    assert "pcli" in result.output.lower() or "portfolio" in result.output.lower()


def test_new_help():
    result = runner.invoke(app, ["new", "--help"])
    assert result.exit_code == 0


def test_validate_help():
    result = runner.invoke(app, ["validate", "--help"])
    assert result.exit_code == 0


def test_ls_help():
    result = runner.invoke(app, ["ls", "--help"])
    assert result.exit_code == 0


def test_show_help():
    result = runner.invoke(app, ["show", "--help"])
    assert result.exit_code == 0


def test_build_resume_help():
    result = runner.invoke(app, ["build", "resume", "--help"])
    assert result.exit_code == 0


def test_build_portfolio_help():
    result = runner.invoke(app, ["build", "portfolio", "--help"])
    assert result.exit_code == 0


# ─── new ───────────────────────────────────────────────────────────────────


def test_new_creates_file(card_repo, monkeypatch):
    import scripts.pcli as pcli_mod

    monkeypatch.setattr(pcli_mod, "REPO_ROOT", card_repo)

    result = runner.invoke(
        app,
        [
            "new",
            "hackathon",
            "my-test-card",
            "--title",
            "My Test Card",
            "--start",
            "2026-06-01",
        ],
    )
    assert result.exit_code == 0
    expected = card_repo / "cards" / "2026-06-my-test-card.mdx"
    assert expected.exists()
    content = expected.read_text(encoding="utf-8")
    assert "id: my-test-card" in content
    assert "type: hackathon" in content


def test_new_invalid_type(card_repo, monkeypatch):
    import scripts.pcli as pcli_mod

    monkeypatch.setattr(pcli_mod, "REPO_ROOT", card_repo)
    result = runner.invoke(app, ["new", "invalidtype", "slug"])
    assert result.exit_code != 0


def test_new_invalid_slug(card_repo, monkeypatch):
    import scripts.pcli as pcli_mod

    monkeypatch.setattr(pcli_mod, "REPO_ROOT", card_repo)
    result = runner.invoke(app, ["new", "hackathon", "Invalid Slug"])
    assert result.exit_code != 0


def test_new_duplicate(card_repo, monkeypatch):
    import scripts.pcli as pcli_mod

    monkeypatch.setattr(pcli_mod, "REPO_ROOT", card_repo)
    # create once
    runner.invoke(app, ["new", "hackathon", "dup-card", "--start", "2026-06-01"])
    # try again
    result = runner.invoke(app, ["new", "hackathon", "dup-card", "--start", "2026-06-01"])
    assert result.exit_code != 0


# ─── validate ──────────────────────────────────────────────────────────────


def test_validate_ok(card_repo, monkeypatch):
    import scripts.pcli as pcli_mod

    monkeypatch.setattr(pcli_mod, "REPO_ROOT", card_repo)
    result = runner.invoke(app, ["validate"])
    assert result.exit_code == 0


def test_validate_error(card_repo, monkeypatch):
    import scripts.pcli as pcli_mod

    monkeypatch.setattr(pcli_mod, "REPO_ROOT", card_repo)
    # write a bad card
    bad = textwrap.dedent("""\
        ---
        id: bad card id
        title: Bad
        type: hackathon
        period:
          start: 2026-05-01
        summary: "ok"
        ---
    """)
    (card_repo / "cards" / "2026-05-bad-card-id.mdx").write_text(bad, encoding="utf-8")
    result = runner.invoke(app, ["validate"])
    assert result.exit_code == 1


def test_validate_slug_broken_card(card_repo, monkeypatch):
    """validate <slug> must report errors for a card that failed parsing (ISSUE-3)."""
    import scripts.pcli as pcli_mod

    monkeypatch.setattr(pcli_mod, "REPO_ROOT", card_repo)
    broken = textwrap.dedent("""\
        ---
        id: broken-card
        title: Broken
        type: hackathon
        period:
          start: 2026-05-01
        summary: "x" * 300
        ---
    """)
    (card_repo / "cards" / "2026-05-broken-card.mdx").write_text(broken, encoding="utf-8")
    # Target the broken card by id slug
    result = runner.invoke(app, ["validate", "broken-card"])
    assert result.exit_code == 1


def test_show_by_filename_stem(card_repo, monkeypatch):
    """show must resolve a card by full filename stem (ISSUE-2 regression)."""
    import scripts.pcli as pcli_mod

    monkeypatch.setattr(pcli_mod, "REPO_ROOT", card_repo)
    result = runner.invoke(app, ["show", "2026-05-sample-card"])
    assert result.exit_code == 0
    assert "Sample Card" in result.output


def test_validate_slug_no_false_positive(card_repo, monkeypatch):
    """validate <slug> must NOT report errors for a different card sharing slug suffix."""
    import scripts.pcli as pcli_mod

    monkeypatch.setattr(pcli_mod, "REPO_ROOT", card_repo)
    # "other-broken-card" shares the "-card" suffix with "broken-card"
    broken = textwrap.dedent("""\
        ---
        id: broken-card
        title: Broken
        type: hackathon
        period:
          start: 2026-05-01
        summary: "x" * 300
        ---
    """)
    other_broken = textwrap.dedent("""\
        ---
        id: other-broken-card
        title: Other Broken
        type: hackathon
        period:
          start: 2026-04-01
        summary: "x" * 300
        ---
    """)
    (card_repo / "cards" / "2026-05-broken-card.mdx").write_text(broken, encoding="utf-8")
    (card_repo / "cards" / "2026-04-other-broken-card.mdx").write_text(
        other_broken, encoding="utf-8"
    )
    # validate "broken-card" must only report broken-card, NOT other-broken-card
    result = runner.invoke(app, ["validate", "broken-card"])
    assert result.exit_code == 1
    assert "other-broken-card" not in result.output


# ─── ls ────────────────────────────────────────────────────────────────────


def test_ls_shows_cards(card_repo, monkeypatch):
    import scripts.pcli as pcli_mod

    monkeypatch.setattr(pcli_mod, "REPO_ROOT", card_repo)
    result = runner.invoke(app, ["ls"])
    assert result.exit_code == 0
    assert "sample-card" in result.output


def test_ls_type_filter(card_repo, monkeypatch):
    import scripts.pcli as pcli_mod

    monkeypatch.setattr(pcli_mod, "REPO_ROOT", card_repo)
    result = runner.invoke(app, ["ls", "--type", "talk"])
    assert result.exit_code == 0
    assert "No cards" in result.output


# ─── show ──────────────────────────────────────────────────────────────────


def test_show_existing(card_repo, monkeypatch):
    import scripts.pcli as pcli_mod

    monkeypatch.setattr(pcli_mod, "REPO_ROOT", card_repo)
    result = runner.invoke(app, ["show", "sample-card"])
    assert result.exit_code == 0
    assert "Sample Card" in result.output


def test_show_missing(card_repo, monkeypatch):
    import scripts.pcli as pcli_mod

    monkeypatch.setattr(pcli_mod, "REPO_ROOT", card_repo)
    result = runner.invoke(app, ["show", "nonexistent-card"])
    assert result.exit_code == 1


# ─── build resume --dry-run ────────────────────────────────────────────────


def test_build_resume_dry_run(card_repo, monkeypatch):
    import scripts.pcli as pcli_mod

    monkeypatch.setattr(pcli_mod, "REPO_ROOT", card_repo)
    result = runner.invoke(app, ["build", "resume", "--dry-run"])
    assert result.exit_code == 0
    assert "sample-card" in result.output
    assert "Dry run" in result.output


# ─── build portfolio ───────────────────────────────────────────────────────


def test_build_portfolio_dry_run(card_repo, monkeypatch):
    import scripts.pcli as pcli_mod

    monkeypatch.setattr(pcli_mod, "REPO_ROOT", card_repo)
    result = runner.invoke(app, ["build", "portfolio", "--dry-run"])
    assert result.exit_code == 0
    assert "sample-card" in result.output
    assert "portfolio" in result.output


def test_build_portfolio_creates_pptx(card_repo, monkeypatch):
    import scripts.pcli as pcli_mod

    monkeypatch.setattr(pcli_mod, "REPO_ROOT", card_repo)
    out = card_repo / "output" / "test-portfolio.pptx"
    result = runner.invoke(app, ["build", "portfolio", "--out", str(out)])
    assert result.exit_code == 0
    assert out.exists()

    from pptx import Presentation

    prs = Presentation(str(out))
    assert len(prs.slides) == 4  # cover, TOC, one card, closing
    text = "\n".join(
        shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")
    )
    assert "Sample Card" in text


def test_build_portfolio_missing_visual_does_not_fail(card_repo, monkeypatch):
    """ISSUE-1: portfolio build succeeds when a card references a non-existent visual."""
    import scripts.pcli as pcli_mod

    monkeypatch.setattr(pcli_mod, "REPO_ROOT", card_repo)
    card_with_visual = textwrap.dedent("""\
        ---
        id: visual-card
        title: Visual Card
        type: project
        period:
          start: 2026-01-01
        status: live
        summary: "Card that has a visual path pointing to a missing file."
        narrative: "Long enough narrative for live status check here yes indeed ok indeed yes."
        evidence:
          - type: repo
            url: https://github.com/example/visual
        visuals:
          - path: assets/nonexistent.png
            role: hero
        ---
    """)
    (card_repo / "cards" / "2026-01-visual-card.mdx").write_text(card_with_visual, encoding="utf-8")
    out = card_repo / "output" / "visual-test.pptx"
    result = runner.invoke(app, ["build", "portfolio", "--out", str(out)])
    assert result.exit_code == 0
    assert out.exists()


def test_build_portfolio_unsupported_layout_exits_nonzero(card_repo, monkeypatch):
    """ISSUE-2: unsupported --layout exits non-zero even with --dry-run."""
    import scripts.pcli as pcli_mod

    monkeypatch.setattr(pcli_mod, "REPO_ROOT", card_repo)
    result = runner.invoke(app, ["build", "portfolio", "--layout", "invalid-layout", "--dry-run"])
    assert result.exit_code != 0


def test_build_portfolio_default_output_in_portfolios_subdir(card_repo, monkeypatch):
    """ISSUE-3: default portfolio output goes to output/portfolios/, not output/ root."""
    import scripts.pcli as pcli_mod

    monkeypatch.setattr(pcli_mod, "REPO_ROOT", card_repo)
    result = runner.invoke(app, ["build", "portfolio"])
    assert result.exit_code == 0
    portfolios_dir = card_repo / "output" / "portfolios"
    assert portfolios_dir.exists()
    assert len(list(portfolios_dir.glob("*.pptx"))) == 1


def test_build_portfolio_explicit_cards_bypass_filters(card_repo, monkeypatch):
    import scripts.pcli as pcli_mod

    monkeypatch.setattr(pcli_mod, "REPO_ROOT", card_repo)
    archived = textwrap.dedent("""\
        ---
        id: archived-card
        title: Archived Card
        type: project
        period:
          start: 2024-01-01
        status: archived
        summary: "Archived but explicitly selected."
        narrative: "This archived narrative is intentionally long enough for validation."
        evidence:
          - type: repo
            url: https://github.com/example/archived
        ---
    """)
    (card_repo / "cards" / "2024-01-archived-card.mdx").write_text(archived, encoding="utf-8")
    out = card_repo / "output" / "explicit.pptx"
    result = runner.invoke(
        app, ["build", "portfolio", "--cards", "archived-card", "--out", str(out)]
    )
    assert result.exit_code == 0

    from pptx import Presentation

    prs = Presentation(str(out))
    text = "\n".join(
        shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")
    )
    assert "Archived Card" in text


# ─── llm commands ─────────────────────────────────────────────────────────


def test_llm_tailor_help():
    result = runner.invoke(app, ["llm", "tailor", "--help"])
    assert result.exit_code == 0


def test_llm_suggest_help():
    result = runner.invoke(app, ["llm", "suggest", "--help"])
    assert result.exit_code == 0


def test_build_resume_dry_run_with_jd(card_repo, monkeypatch, tmp_path):
    """--jd + --dry-run shows LLM intent without calling the LLM."""
    import scripts.pcli as pcli_mod

    monkeypatch.setattr(pcli_mod, "REPO_ROOT", card_repo)
    jd_file = tmp_path / "jd.txt"
    jd_file.write_text("Senior engineer needed.", encoding="utf-8")
    result = runner.invoke(
        app, ["build", "resume", "--jd", str(jd_file), "--tone", "formal", "--dry-run"]
    )
    assert result.exit_code == 0
    assert "jd=" in result.output or "Dry run" in result.output


def test_llm_tailor_with_fake_client(card_repo, monkeypatch, tmp_path):
    """llm tailor prints score + rewrite using a fake client."""
    import json

    import scripts.pcli as pcli_mod

    monkeypatch.setattr(pcli_mod, "REPO_ROOT", card_repo)

    class _FakeContent:
        def __init__(self, text):
            self.text = text

    class _FakeMsg:
        def __init__(self, text):
            self.content = [_FakeContent(text)]

    responses = [
        json.dumps({"score": 0.9, "reason": "great match"}),
        "Rewritten summary for formal tone.",
    ]
    resp_iter = iter(responses)

    class _FakeMsgs:
        def create(self, **kwargs):
            return _FakeMsg(next(resp_iter))

    class _FakeClient:
        messages = _FakeMsgs()

    import scripts.llm as llm_mod

    monkeypatch.setattr(llm_mod, "_build_client", lambda *a, **k: _FakeClient())
    monkeypatch.setattr(pcli_mod, "REPO_ROOT", card_repo)

    jd_file = tmp_path / "jd.txt"
    jd_file.write_text("Senior engineer.", encoding="utf-8")
    result = runner.invoke(
        app,
        [
            "llm",
            "tailor",
            "--cards",
            "sample-card",
            "--jd",
            str(jd_file),
            "--no-cache",
        ],
    )
    assert result.exit_code == 0
    assert "sample-card" in result.output
    assert "0.90" in result.output or "0.9" in result.output


def test_llm_suggest_with_fake_client(card_repo, monkeypatch, tmp_path):
    """llm suggest prints YAML frontmatter from fake LLM response."""
    import json

    import scripts.pcli as pcli_mod

    monkeypatch.setattr(pcli_mod, "REPO_ROOT", card_repo)

    fake_response = json.dumps(
        {
            "title": "Bridge Contract",
            "type": "project",
            "summary": "Built a cross-chain bridge.",
            "tags_domain": ["web3"],
            "tags_skill": ["solidity"],
        }
    )

    class _FakeContent:
        def __init__(self, text):
            self.text = text

    class _FakeMsg:
        def __init__(self, text):
            self.content = [_FakeContent(text)]

    class _FakeMsgs:
        def create(self, **kwargs):
            return _FakeMsg(fake_response)

    class _FakeClient:
        messages = _FakeMsgs()

    import scripts.llm as llm_mod

    monkeypatch.setattr(llm_mod, "_build_client", lambda *a, **k: _FakeClient())

    notes_file = tmp_path / "notes.txt"
    notes_file.write_text("Built a cross-chain bridge between ETH and Solana.", encoding="utf-8")
    result = runner.invoke(app, ["llm", "suggest", "--from", str(notes_file), "--no-cache"])
    assert result.exit_code == 0
    assert "Bridge Contract" in result.output


def test_llm_tailor_missing_api_key(card_repo, monkeypatch, tmp_path):
    """llm tailor exits 2 when no API key and no cache hit."""
    import scripts.pcli as pcli_mod

    monkeypatch.setattr(pcli_mod, "REPO_ROOT", card_repo)
    monkeypatch.delenv("AI_PROVIDER", raising=False)
    monkeypatch.delenv("ANTHROPIC_API_KEY", raising=False)
    monkeypatch.delenv("AI_API_KEY", raising=False)
    monkeypatch.delenv("GOOGLE_API_KEY", raising=False)
    monkeypatch.delenv("GEMINI_API_KEY", raising=False)

    jd_file = tmp_path / "jd.txt"
    jd_file.write_text("Engineer.", encoding="utf-8")
    result = runner.invoke(
        app,
        ["llm", "tailor", "--cards", "sample-card", "--jd", str(jd_file), "--no-cache"],
    )
    assert result.exit_code == 2


def test_build_resume_llm_meta_persisted_in_resume_data_json(card_repo, monkeypatch, tmp_path):
    """After LLM resume build, .build/resume-data.json must contain meta.llm."""
    import json

    import scripts.llm as llm_mod
    import scripts.pcli as pcli_mod

    monkeypatch.setattr(pcli_mod, "REPO_ROOT", card_repo)

    # Stub default.typ so template existence check passes
    (card_repo / "templates" / "resume").mkdir(parents=True, exist_ok=True)
    (card_repo / "templates" / "resume" / "default.typ").write_text(
        "#let data = none", encoding="utf-8"
    )

    # Fake LLM client: score response then rewrite response
    score_resp = json.dumps({"score": 0.85, "reason": "good match"})
    rewrite_resp = "LLM-rewritten summary."

    class _FC:
        def __init__(self, text):
            self.text = text

    class _FM:
        def __init__(self, text):
            self.content = [_FC(text)]

    responses = iter([score_resp, rewrite_resp])

    class _FakeClient:
        class messages:
            @staticmethod
            def create(**kwargs):
                return _FM(next(responses))

    monkeypatch.setattr(llm_mod, "_build_client", lambda *a, **k: _FakeClient())

    # Fake typst so build_resume doesn't fail on missing binary
    import shutil

    monkeypatch.setattr(shutil, "which", lambda _: "/fake/typst")

    import subprocess

    monkeypatch.setattr(
        subprocess,
        "run",
        lambda *a, **kw: type("R", (), {"stdout": "", "stderr": "", "returncode": 0})(),
    )

    jd_file = tmp_path / "jd.txt"
    jd_file.write_text("Blockchain engineer needed.", encoding="utf-8")
    out = card_repo / "output" / "llm-test.pdf"

    result = runner.invoke(
        app,
        [
            "build",
            "resume",
            "--jd",
            str(jd_file),
            "--tone",
            "formal",
            "--out",
            str(out),
            "--no-cache",
        ],
    )
    assert result.exit_code == 0, result.output

    data = json.loads((card_repo / ".build" / "resume-data.json").read_text(encoding="utf-8"))
    assert "llm" in data["meta"]
    assert "scores" in data["meta"]["llm"]
    assert "rewrites" in data["meta"]["llm"]
