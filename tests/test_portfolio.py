"""Tests for the portfolio PPTX renderer."""

from __future__ import annotations

from datetime import date

from pptx import Presentation

from scripts.card import Card
from scripts.render_portfolio import PortfolioOptions
from templates.portfolio.default import render


def _card(**overrides) -> Card:
    data = {
        "id": "visual-card",
        "title": "Visual Card",
        "type": "project",
        "period": {"start": date(2026, 5, 1)},
        "status": "live",
        "tags": {"domain": ["web3"], "skill": [], "audience": ["web3"]},
        "summary": "A concise summary for the portfolio slide.",
        "narrative": (
            "NarrativeOnlyMarker explains context, approach, and outcome in a longer "
            "paragraph intended only for the portfolio narrative area."
        ),
        "evidence": [{"type": "repo", "url": "https://github.com/example/repo"}],
    }
    data.update(overrides)
    return Card.model_validate(data)


def _deck_text(path) -> str:
    prs = Presentation(str(path))
    return "\n".join(
        shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")
    )


def test_missing_visual_renders_placeholder(tmp_path):
    card = _card(visuals=[{"path": "assets/missing.png", "role": "hero"}])
    out = tmp_path / "missing-visual.pptx"
    render(
        cards=[card],
        profile={"basics": {"name": "Test User", "email": "test@example.com"}},
        options=PortfolioOptions(repo_root=tmp_path),
        output_path=out,
    )

    assert out.exists()
    text = _deck_text(out)
    assert "Missing visual: assets/missing.png" in text
    assert "Visual Card" in text


def test_no_narrative_excludes_narrative_text(tmp_path):
    card = _card()
    out = tmp_path / "no-narrative.pptx"
    render(
        cards=[card],
        profile={"basics": {"name": "Test User", "email": "test@example.com"}},
        options=PortfolioOptions(include_narrative=False, repo_root=tmp_path),
        output_path=out,
    )

    text = _deck_text(out)
    assert "Visual Card" in text
    assert "NarrativeOnlyMarker" not in text
