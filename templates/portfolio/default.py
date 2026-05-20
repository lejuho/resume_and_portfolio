"""Default python-pptx portfolio renderer."""

from __future__ import annotations

import re
from datetime import date
from pathlib import Path
from textwrap import shorten

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from rich.console import Console

from scripts.card import Card
from scripts.render_portfolio import PortfolioOptions

console = Console()


THEMES = {
    "web3": {"bg": "101827", "fg": "F8FAFC", "muted": "CBD5E1", "accent": "00FFA3"},
    "regulatory": {"bg": "F8F5EC", "fg": "172554", "muted": "475569", "accent": "1D4ED8"},
    "korean-market": {"bg": "F8F5EC", "fg": "172554", "muted": "475569", "accent": "1D4ED8"},
    "default": {"bg": "F8FAFC", "fg": "111827", "muted": "64748B", "accent": "2563EB"},
}


def render(
    cards: list[Card],
    profile: dict,
    options: PortfolioOptions,
    output_path: Path,
) -> None:
    layout = options.layout
    if layout == "grouped-by-type":
        _render_grouped(cards, profile, options, output_path)
    elif layout == "timeline":
        _render_timeline(cards, profile, options, output_path)
    else:
        _render_default(cards, profile, options, output_path)


def _render_default(
    cards: list[Card],
    profile: dict,
    options: PortfolioOptions,
    output_path: Path,
) -> None:
    prs = _new_presentation()
    theme = _theme_for(cards)
    _add_cover(prs, cards, profile, options, theme)
    _add_toc(prs, cards, theme)
    for index, card in enumerate(cards, start=1):
        _add_card_slide(prs, card, index, options, theme)
    _add_closing(prs, profile, theme)
    prs.save(output_path)


def _render_grouped(
    cards: list[Card],
    profile: dict,
    options: PortfolioOptions,
    output_path: Path,
) -> None:
    """grouped-by-type: cover → TOC → (type divider → cards per type) → closing."""
    from itertools import groupby

    prs = _new_presentation()
    theme = _theme_for(cards)
    _add_cover(prs, cards, profile, options, theme)
    _add_toc(prs, cards, theme)

    sorted_cards = sorted(cards, key=lambda c: c.type)
    global_index = 1
    for card_type, group in groupby(sorted_cards, key=lambda c: c.type):
        group_list = list(group)
        _add_type_divider(prs, card_type, len(group_list), theme)
        for card in group_list:
            _add_card_slide(prs, card, global_index, options, theme)
            global_index += 1

    _add_closing(prs, profile, theme)
    prs.save(output_path)


def _render_timeline(
    cards: list[Card],
    profile: dict,
    options: PortfolioOptions,
    output_path: Path,
) -> None:
    """timeline: cover → timeline overview → highlight slides (live + metrics) → closing."""
    prs = _new_presentation()
    theme = _theme_for(cards)
    _add_cover(prs, cards, profile, options, theme)
    _add_timeline_overview(prs, cards, theme)

    highlights = [c for c in cards if c.status == "live" and c.metrics]
    highlights_sorted = sorted(highlights, key=lambda c: c.period.start, reverse=True)
    for index, card in enumerate(highlights_sorted, start=1):
        _add_card_slide(prs, card, index, options, theme)

    _add_closing(prs, profile, theme)
    prs.save(output_path)


def _new_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def _theme_for(cards: list[Card]) -> dict[str, str]:
    audiences = cards[0].tags.audience if cards else []
    for audience in audiences:
        if audience in THEMES:
            return THEMES[audience]
    return THEMES["default"]


def _rgb(value: str) -> RGBColor:
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def _fill_bg(slide, theme: dict[str, str]) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _rgb(theme["bg"])


def _textbox(slide, x, y, w, h, text: str, size=18, color="111827", bold=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = _rgb(color)
    return box


def _add_cover(prs, cards: list[Card], profile: dict, options: PortfolioOptions, theme):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_bg(slide, theme)
    basics = profile.get("basics", {})
    title = options.cover_title or basics.get("label") or "Portfolio"
    subtitle = options.cover_subtitle or f"{date.today().isoformat()} · {len(cards)} selected cards"

    _textbox(slide, 0.9, 1.25, 10.8, 0.8, title, size=32, color=theme["fg"], bold=True)
    _textbox(slide, 0.95, 2.2, 8.8, 0.45, subtitle, size=17, color=theme["muted"])
    _textbox(slide, 0.95, 5.65, 7.0, 0.4, basics.get("name", ""), size=18, color=theme["fg"])
    _accent_bar(slide, theme)


def _add_toc(prs, cards: list[Card], theme):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_bg(slide, theme)
    _textbox(
        slide, 0.65, 0.55, 5.0, 0.45, "Table of Contents", size=24, color=theme["fg"], bold=True
    )
    for i, card in enumerate(cards, start=1):
        y = 1.2 + (i - 1) * 0.52
        _textbox(slide, 0.85, y, 0.45, 0.28, f"{i:02}", size=11, color=theme["accent"], bold=True)
        _textbox(slide, 1.35, y - 0.04, 8.8, 0.32, card.title, size=15, color=theme["fg"])
        _textbox(slide, 10.4, y, 1.5, 0.25, card.type, size=10, color=theme["muted"])
    _accent_bar(slide, theme)


def _add_card_slide(prs, card: Card, index: int, options: PortfolioOptions, theme):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_bg(slide, theme)

    _add_visual(slide, card, options, theme)
    _textbox(slide, 6.9, 0.55, 5.7, 0.65, card.title, size=23, color=theme["fg"], bold=True)
    _textbox(
        slide,
        6.95,
        1.25,
        4.0,
        0.28,
        f"{card.type} · {_period(card)}",
        size=11,
        color=theme["muted"],
    )
    _textbox(slide, 6.95, 1.75, 5.2, 0.75, card.summary.strip(), size=13, color=theme["fg"])

    y = 2.75
    if card.metrics:
        _textbox(slide, 6.95, y, 1.4, 0.25, "Metrics", size=10, color=theme["accent"], bold=True)
        y += 0.3
        for metric in card.metrics[:3]:
            _textbox(slide, 7.05, y, 5.4, 0.27, f"- {metric}", size=10.5, color=theme["fg"])
            y += 0.3

    if options.include_narrative and card.effective_narrative:
        narrative = _clean_narrative(card.effective_narrative)
        _textbox(
            slide,
            6.95,
            y + 0.12,
            5.4,
            1.15,
            shorten(narrative, width=360),
            size=10.5,
            color=theme["muted"],
        )

    tags = card.tags.domain + card.tags.skill + card.tags.audience
    _textbox(slide, 6.95, 5.75, 5.4, 0.3, " · ".join(tags[:7]), size=9.5, color=theme["accent"])
    evidence = " | ".join(e.url for e in card.evidence[:2])
    _textbox(slide, 6.95, 6.28, 5.4, 0.3, evidence, size=8.5, color=theme["muted"])
    _textbox(slide, 12.35, 6.85, 0.4, 0.2, f"{index}", size=8, color=theme["muted"])


def _add_visual(slide, card: Card, options: PortfolioOptions, theme):
    x, y, w, h = 0.65, 0.65, 5.75, 6.15
    visual = next(
        (v for v in card.visuals if v.role == "hero"), card.visuals[0] if card.visuals else None
    )
    visual_path = options.repo_root / visual.path if visual else None

    if visual_path and visual_path.exists():
        slide.shapes.add_picture(
            str(visual_path), Inches(x), Inches(y), width=Inches(w), height=Inches(h)
        )
        return

    rect = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    rect.fill.solid()
    rect.fill.fore_color.rgb = _rgb("E2E8F0")
    rect.line.color.rgb = _rgb(theme["accent"])
    label = "No visual provided"
    if visual_path:
        label = f"Missing visual: {visual.path}"
        console.print(f"[yellow]WARN[/yellow] {card.id}: missing visual {visual.path}")
    box = _textbox(slide, x + 0.35, y + 2.7, w - 0.7, 0.4, label, size=14, color="334155")
    box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


def _add_type_divider(prs, card_type: str, count: int, theme):
    """Section divider slide for grouped-by-type layout."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_bg(slide, theme)
    _textbox(
        slide, 1.2, 2.8, 8.0, 0.8, card_type.upper(), size=36, color=theme["accent"], bold=True
    )
    _textbox(slide, 1.2, 3.75, 5.0, 0.4, f"{count} card(s)", size=16, color=theme["muted"])
    _accent_bar(slide, theme)


def _add_timeline_overview(prs, cards: list[Card], theme):
    """Timeline overview slide listing all cards chronologically."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_bg(slide, theme)
    _textbox(slide, 0.65, 0.45, 6.0, 0.5, "Timeline", size=24, color=theme["fg"], bold=True)
    sorted_cards = sorted(cards, key=lambda c: c.period.start)
    for i, card in enumerate(sorted_cards[:10]):
        y = 1.1 + i * 0.56
        start = card.period.start.isoformat()
        _textbox(slide, 0.85, y, 1.6, 0.3, start, size=10, color=theme["muted"])
        _textbox(slide, 2.55, y, 7.5, 0.32, card.title, size=12, color=theme["fg"])
        _textbox(slide, 10.2, y, 2.5, 0.25, card.type, size=9, color=theme["accent"])
    _accent_bar(slide, theme)


def _add_closing(prs, profile: dict, theme):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_bg(slide, theme)
    basics = profile.get("basics", {})
    _textbox(slide, 0.9, 1.5, 8.0, 0.6, "Thank you", size=34, color=theme["fg"], bold=True)
    contact = " · ".join(
        str(basics[k]) for k in ("email", "phone", "location", "url") if basics.get(k)
    )
    _textbox(slide, 0.95, 2.45, 9.5, 0.4, contact, size=15, color=theme["muted"])
    _textbox(slide, 0.95, 5.7, 7.0, 0.35, basics.get("name", ""), size=18, color=theme["fg"])
    _accent_bar(slide, theme)


def _accent_bar(slide, theme):
    bar = slide.shapes.add_shape(1, Inches(0), Inches(7.18), Inches(13.333), Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = _rgb(theme["accent"])
    bar.line.fill.background()


def _period(card: Card) -> str:
    end = card.period.end.isoformat() if card.period.end else "present"
    return f"{card.period.start.isoformat()} - {end}"


def _clean_narrative(text: str) -> str:
    text = re.sub(r"^#+\s*", "", text, flags=re.MULTILINE)
    return " ".join(text.split())
